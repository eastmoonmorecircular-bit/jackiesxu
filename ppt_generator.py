"""把 deep_summary 渲染成 .pptx 文件（含 PDF 提取的关键图）"""
from __future__ import annotations
from pathlib import Path
from typing import Dict, List, Optional

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image

from utils import setup_logger

logger = setup_logger(__name__)

OUT_DIR = Path("data/ppts")
OUT_DIR.mkdir(parents=True, exist_ok=True)

# 颜色主题
COLOR_PRIMARY = RGBColor(0x1F, 0x4E, 0x79)   # 深蓝
COLOR_ACCENT = RGBColor(0xE7, 0x6F, 0x00)    # 橙
COLOR_TEXT = RGBColor(0x33, 0x33, 0x33)
COLOR_LIGHT = RGBColor(0x66, 0x66, 0x66)
COLOR_BG_LIGHT = RGBColor(0xF5, 0xF7, 0xFA)


# =================== 通用 helper ===================

def _add_title_bar(slide, title: str, subtitle: str = ""):
    """每页顶部统一标题带"""
    # 顶部色块
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.33), Inches(0.7))
    bar.fill.solid()
    bar.fill.fore_color.rgb = COLOR_PRIMARY
    bar.line.fill.background()

    # 标题
    tx = slide.shapes.add_textbox(Inches(0.4), Inches(0.1), Inches(12.5), Inches(0.6))
    tf = tx.text_frame
    tf.margin_left = tf.margin_right = 0
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.size = Pt(24)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    if subtitle:
        p2 = tf.add_paragraph()
        run2 = p2.add_run()
        run2.text = subtitle
        run2.font.size = Pt(11)
        run2.font.color.rgb = RGBColor(0xCC, 0xDD, 0xEE)


def _add_bullets(slide, items: List[str], left=0.5, top=1.0, width=12.3, height=5.8,
                 size=16, accent_color=None):
    """添加 bullet 列表"""
    tx = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tx.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if not item:
            continue
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(8)
        # bullet 符号
        marker = p.add_run()
        marker.text = "▸ "
        marker.font.color.rgb = accent_color or COLOR_ACCENT
        marker.font.bold = True
        marker.font.size = Pt(size)
        # 内容
        run = p.add_run()
        run.text = str(item)
        run.font.size = Pt(size)
        run.font.color.rgb = COLOR_TEXT


def _add_section_text(slide, label: str, value: str,
                      left=0.5, top=1.0, width=12.3, height=2.0, size=14):
    """添加带 label 的段落（label 加粗 + 主色，value 黑色正文）"""
    tx = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tx.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    if label:
        rl = p.add_run()
        rl.text = f"{label}： "
        rl.font.size = Pt(size + 1)
        rl.font.bold = True
        rl.font.color.rgb = COLOR_PRIMARY
    rv = p.add_run()
    rv.text = str(value or "")
    rv.font.size = Pt(size)
    rv.font.color.rgb = COLOR_TEXT


def _add_image(slide, img_path: str, left, top, max_width=8.0, max_height=4.5):
    """智能插图（自动等比缩放）"""
    try:
        with Image.open(img_path) as im:
            w_px, h_px = im.size
        ratio = w_px / max(h_px, 1)
        target_w = max_width
        target_h = target_w / ratio
        if target_h > max_height:
            target_h = max_height
            target_w = target_h * ratio
        slide.shapes.add_picture(img_path, Inches(left), Inches(top),
                                 width=Inches(target_w), height=Inches(target_h))
    except Exception as e:
        logger.warning(f"插图失败 {img_path}: {e}")


def _add_footer(slide, text: str):
    tx = slide.shapes.add_textbox(Inches(0.4), Inches(7.0), Inches(12.5), Inches(0.3))
    p = tx.text_frame.paragraphs[0]
    p.alignment = 2  # right
    run = p.add_run()
    run.text = text
    run.font.size = Pt(9)
    run.font.color.rgb = COLOR_LIGHT


# =================== 各章节页 ===================

def _slide_cover(prs, paper: dict):
    s = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    # 大色块背景
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.33), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = COLOR_PRIMARY
    bg.line.fill.background()

    # 标签
    tag = s.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(11.7), Inches(0.5))
    p = tag.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = "推荐系统论文深度解读"
    run.font.size = Pt(18)
    run.font.color.rgb = COLOR_ACCENT

    # 标题
    title_box = s.shapes.add_textbox(Inches(0.8), Inches(2.2), Inches(11.7), Inches(2.5))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = paper["title"]
    run.font.size = Pt(32)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    # 元信息
    meta = s.shapes.add_textbox(Inches(0.8), Inches(5.5), Inches(11.7), Inches(1.5))
    tf = meta.text_frame
    for line in [
        f"作者：{paper.get('authors', '')[:200]}",
        f"arXiv：{paper.get('arxiv_id', '')}  |  发布：{paper.get('published', '')}",
        f"方向：{' · '.join(paper.get('directions') or ['其他'])}",
    ]:
        p = tf.add_paragraph() if tf.paragraphs[0].text else tf.paragraphs[0]
        run = p.add_run()
        run.text = line
        run.font.size = Pt(13)
        run.font.color.rgb = RGBColor(0xCC, 0xDD, 0xEE)


def _slide_overview(prs, paper: dict):
    """概览页：浅解读的 problem / method / innovation"""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _add_title_bar(s, "📌 一页速览 TL;DR", paper.get("arxiv_id", ""))
    summary = paper.get("summary") or {}

    items = []
    if summary.get("problem"):
        items.append(f"【问题】{summary['problem']}")
    if summary.get("method"):
        items.append(f"【方法】{summary['method']}")
    if summary.get("innovation"):
        for innov in summary["innovation"][:3]:
            items.append(f"【创新】{innov}")
    if summary.get("experiment"):
        items.append(f"【实验】{summary['experiment']}")
    score = summary.get("value_score")
    prio = summary.get("reading_priority")
    if score or prio:
        items.append(f"【价值评分】{'⭐'*int(score or 0)}    【阅读优先级】{prio or '-'}")
    _add_bullets(s, items, top=1.1, size=15)


def _slide_motivation(prs, deep: dict):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _add_title_bar(s, "🎯 Motivation 动机解读")
    m = deep.get("motivation") or {}
    if m.get("_parse_failed"):
        _add_bullets(s, [m.get("_raw", "解析失败")], top=1.1)
        return
    _add_section_text(s, "现实痛点", m.get("real_world_pain", ""), top=1.1, height=1.2)
    items = m.get("existing_limitations", []) or []
    _add_section_text(s, "现有方法不足", "", top=2.5, height=0.4, size=14)
    _add_bullets(s, items, top=3.0, height=2.5, size=14)
    _add_section_text(s, "本文切入角度", m.get("this_paper_angle", ""), top=5.5, height=0.9)
    _add_section_text(s, "为什么重要", m.get("why_it_matters", ""), top=6.4, height=0.6, size=12)


def _slide_problem(prs, deep: dict):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _add_title_bar(s, "❓ 问题定义 Problem Formulation")
    p = deep.get("problem") or {}
    if p.get("_parse_failed"):
        _add_bullets(s, [p.get("_raw", "")], top=1.1)
        return
    _add_section_text(s, "任务类型", p.get("task_type", ""), top=1.1, height=0.6)
    _add_section_text(s, "形式化定义", p.get("formal_definition", ""), top=1.8, height=2.0)
    _add_section_text(s, "输入", p.get("input", ""), top=4.0, height=0.8)
    _add_section_text(s, "输出", p.get("output", ""), top=4.9, height=0.8)
    items = p.get("key_constraints", []) or []
    if items:
        _add_section_text(s, "关键约束", "", top=5.8, height=0.4)
        _add_bullets(s, items, top=6.2, height=1.3, size=13)


def _slide_method(prs, deep: dict, image_paths: List[str]):
    """方法详解（多页）：概览 + 模块 + 训练/推理"""
    m = deep.get("method") or {}
    if m.get("_parse_failed"):
        s = prs.slides.add_slide(prs.slide_layouts[6])
        _add_title_bar(s, "🛠 方法解读")
        _add_bullets(s, [m.get("_raw", "")], top=1.1)
        return

    # Page 1: 整体框架 + 关键洞见 + 架构图（如有）
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _add_title_bar(s, "🛠 方法概览 Method Overview")
    _add_section_text(s, "整体框架", m.get("overview", ""), top=1.1, height=2.0)
    insights = m.get("key_insights", []) or []
    if insights:
        _add_section_text(s, "关键洞见", "", top=3.2, height=0.4)
        _add_bullets(s, insights, top=3.7, height=2.0, size=13)
    # 架构图（第一张大图）
    if image_paths:
        _add_image(s, image_paths[0], left=4.5, top=5.7, max_width=4.3, max_height=1.5)

    # Page 2+: 关键模块（每页 2-3 个模块）
    modules = m.get("key_modules", []) or []
    if modules:
        for chunk_start in range(0, len(modules), 3):
            chunk = modules[chunk_start:chunk_start + 3]
            s = prs.slides.add_slide(prs.slide_layouts[6])
            _add_title_bar(s, "🔬 方法细节 Key Modules", f"模块 {chunk_start+1}-{chunk_start+len(chunk)}")
            top = 1.1
            for mod in chunk:
                name = mod.get("name", "")
                role = mod.get("role", "")
                how = mod.get("how", "")
                formula = mod.get("formula", "")
                _add_section_text(s, f"▌ {name}", role, top=top, height=0.6, size=14)
                _add_section_text(s, "  实现", how, top=top + 0.6, height=1.2, size=12)
                if formula:
                    _add_section_text(s, "  关键公式", formula, top=top + 1.7, height=0.5, size=11)
                    top += 2.2
                else:
                    top += 1.85

    # Page 3: 训练与推理
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _add_title_bar(s, "🧪 训练与推理")
    _add_section_text(s, "训练策略", m.get("training_strategy", ""), top=1.1, height=2.5)
    _add_section_text(s, "推理流程", m.get("inference_pipeline", ""), top=3.8, height=2.5)


def _slide_experiment(prs, deep: dict, image_paths: List[str]):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _add_title_bar(s, "📊 实验设置 & 主结果")
    e = deep.get("experiment") or {}
    if e.get("_parse_failed"):
        _add_bullets(s, [e.get("_raw", "")], top=1.1)
        return
    # 数据集
    ds = e.get("datasets", []) or []
    ds_text = " | ".join(f"{d.get('name','')}({d.get('scale','')})" for d in ds)
    _add_section_text(s, "数据集", ds_text, top=1.1, height=0.7, size=13)
    # 指标 / baseline
    _add_section_text(s, "评价指标", " · ".join(e.get("metrics") or []), top=1.9, height=0.7, size=13)
    bls = e.get("baselines", []) or []
    bl_text = " | ".join(f"{b.get('name','')}({b.get('type','')})" for b in bls)
    _add_section_text(s, "对比基线", bl_text, top=2.7, height=1.0, size=12)
    # 主结果
    _add_section_text(s, "主要发现", e.get("main_results", ""), top=3.9, height=1.8, size=14)
    _add_section_text(s, "胜负分析", e.get("win_or_lose", ""), top=5.8, height=1.3, size=12)
    # 实验图（第二张）
    if len(image_paths) > 1:
        _add_image(s, image_paths[1], left=8.5, top=5.5, max_width=4.5, max_height=1.5)


def _slide_ablation(prs, deep: dict):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _add_title_bar(s, "🔍 消融实验逐项分析 Ablation Study")
    a = deep.get("ablation") or {}
    if a.get("_parse_failed"):
        _add_bullets(s, [a.get("_raw", "")], top=1.1)
        return
    abls = a.get("ablations", []) or []
    items = []
    for i, ab in enumerate(abls[:6], 1):
        items.append(
            f"#{i} 去掉 [{ab.get('removed_component','')}] → "
            f"{ab.get('impact','')} | 解读：{ab.get('interpretation','')}"
        )
    _add_bullets(s, items, top=1.1, height=4.0, size=12)
    _add_section_text(s, "超参敏感性", a.get("hyperparameter_analysis", ""), top=5.3, height=0.8, size=12)
    _add_section_text(s, "鲁棒性", a.get("robustness", ""), top=6.2, height=0.7, size=12)


def _slide_comparison(prs, deep: dict):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _add_title_bar(s, "⚖️ 与 SOTA 对比：优缺点 Comparison")
    c = deep.get("comparison") or {}
    if c.get("_parse_failed"):
        _add_bullets(s, [c.get("_raw", "")], top=1.1)
        return
    items = []
    for v in (c.get("vs_sota") or [])[:5]:
        items.append(
            f"vs {v.get('competitor','')} —— ✅ {v.get('advantage','')}  ❌ {v.get('disadvantage','')}"
        )
    _add_bullets(s, items, top=1.1, height=3.5, size=12)
    _add_section_text(s, "复杂度对比", c.get("complexity_compare", ""), top=4.8, height=0.8, size=12)
    _add_section_text(s, "工程友好度", c.get("engineering_friendliness", ""), top=5.7, height=0.7, size=12)
    _add_section_text(s, "工业落地价值", c.get("industry_value", ""), top=6.5, height=0.7, size=12)


def _slide_limitation(prs, deep: dict):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _add_title_bar(s, "⚠️ 局限 & 未来方向 Limitations")
    l = deep.get("limitation") or {}
    if l.get("_parse_failed"):
        _add_bullets(s, [l.get("_raw", "")], top=1.1)
        return
    _add_section_text(s, "局限", "", top=1.1, height=0.4)
    _add_bullets(s, l.get("limitations", []) or [], top=1.5, height=2.0, size=13)
    _add_section_text(s, "未来方向", "", top=3.7, height=0.4)
    _add_bullets(s, l.get("future_directions", []) or [], top=4.1, height=1.5, size=13)
    _add_section_text(s, "开放问题", "", top=5.7, height=0.4)
    _add_bullets(s, l.get("open_questions", []) or [], top=6.1, height=0.9, size=12)
    _add_section_text(s, "📖 阅读建议", l.get("reading_advice", ""), top=7.0, height=0.4, size=11)


def _slide_appendix_images(prs, image_paths: List[str]):
    """附录页：把剩下的关键图都展示一下"""
    if len(image_paths) <= 2:
        return
    for i in range(2, len(image_paths)):
        s = prs.slides.add_slide(prs.slide_layouts[6])
        _add_title_bar(s, f"📎 附图 Figure {i}", "来自原文")
        _add_image(s, image_paths[i], left=2.0, top=1.2, max_width=9.5, max_height=5.5)


# =================== 主接口 ===================

def generate_ppt(paper: dict, out_path: Optional[Path] = None) -> Path:
    """根据 paper（含 deep_summary）生成 PPT 文件"""
    deep = paper.get("deep_summary") or {}
    if not deep:
        raise ValueError("paper 缺少 deep_summary，请先运行深度解读")

    image_paths = (deep.get("_meta") or {}).get("image_paths", [])

    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    _slide_cover(prs, paper)
    _slide_overview(prs, paper)
    _slide_motivation(prs, deep)
    _slide_problem(prs, deep)
    _slide_method(prs, deep, image_paths)
    _slide_experiment(prs, deep, image_paths)
    _slide_ablation(prs, deep)
    _slide_comparison(prs, deep)
    _slide_limitation(prs, deep)
    _slide_appendix_images(prs, image_paths)

    if out_path is None:
        safe_id = paper["arxiv_id"].replace("/", "_")
        out_path = OUT_DIR / f"{safe_id}.pptx"
    prs.save(out_path)
    logger.info(f"PPT 已生成: {out_path}")
    return out_path
